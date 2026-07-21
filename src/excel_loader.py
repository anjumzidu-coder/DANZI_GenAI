from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
import tempfile
from typing import Dict

import pandas as pd
from pypdf import PdfReader
from pptx import Presentation


@dataclass
class WorkbookData:
    sheets: Dict[str, pd.DataFrame]


def _pptx_to_dataframe(pptx_path: str) -> pd.DataFrame:
    presentation = Presentation(pptx_path)
    rows = []
    for i, slide in enumerate(presentation.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text.strip())
        rows.append({"slide": i, "text": "\n".join([t for t in text_parts if t])})
    return pd.DataFrame(rows)


def _ppt_to_pptx_and_extract(ppt_bytes: bytes) -> pd.DataFrame:
    with tempfile.TemporaryDirectory() as tmp_dir:
        ppt_path = str(Path(tmp_dir) / "upload.ppt")
        pptx_path = str(Path(tmp_dir) / "converted.pptx")

        with open(ppt_path, "wb") as f:
            f.write(ppt_bytes)

        try:
            import win32com.client  # type: ignore
        except ImportError as exc:
            raise RuntimeError(
                "Legacy .ppt requires pywin32 and PowerPoint installed. "
                "Please upload .pptx if conversion is unavailable."
            ) from exc

        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = None
        try:
            presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=False)
            # 24 = ppSaveAsOpenXMLPresentation (.pptx)
            presentation.SaveAs(pptx_path, 24)
        finally:
            if presentation is not None:
                presentation.Close()
            powerpoint.Quit()

        return _pptx_to_dataframe(pptx_path)


def load_workbook(file_bytes: bytes, file_name: str) -> WorkbookData:
    lower = file_name.lower()

    if lower.endswith(".csv"):
        df = pd.read_csv(BytesIO(file_bytes))
        return WorkbookData(sheets={"CSV": df})

    if lower.endswith(".pdf"):
        reader = PdfReader(BytesIO(file_bytes))
        rows = []
        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            rows.append({"page": i, "text": text.strip()})
        df = pd.DataFrame(rows)
        return WorkbookData(sheets={"PDF_Text": df})

    if lower.endswith(".pptx"):
        with tempfile.TemporaryDirectory() as tmp_dir:
            pptx_path = str(Path(tmp_dir) / "upload.pptx")
            with open(pptx_path, "wb") as f:
                f.write(file_bytes)
            df = _pptx_to_dataframe(pptx_path)
        return WorkbookData(sheets={"PPT_Text": df})

    if lower.endswith(".ppt"):
        df = _ppt_to_pptx_and_extract(file_bytes)
        return WorkbookData(sheets={"PPT_Text": df})

    excel = pd.ExcelFile(BytesIO(file_bytes))
    sheets: Dict[str, pd.DataFrame] = {}
    for sheet in excel.sheet_names:
        sheets[sheet] = excel.parse(sheet)
    return WorkbookData(sheets=sheets)
